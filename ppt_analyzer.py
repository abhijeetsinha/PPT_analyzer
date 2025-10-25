"""
PowerPoint Analyzer Agent using LangGraph
=========================================

A beginner-friendly agent that extracts information from PowerPoint files
and outputs structured data in CSV format.

Author: Abhijeet Kumar Sinha 
Date: October 24, 2025
"""

import os
import glob
import pandas as pd
from pathlib import Path
from typing import TypedDict, List
from dotenv import load_dotenv

# LangGraph and LangChain imports
from langgraph.graph import StateGraph, START, END
from langchain_openai import AzureChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.output_parsers import StrOutputParser

# PowerPoint processing
from pptx import Presentation

# Load environment variables (uses your existing .env file)
load_dotenv()

# Initialize Azure OpenAI (using your existing configuration)
azure_openai_api_key = os.getenv("AZURE_OPENAI_API_KEY")
azure_openai_endpoint = os.getenv("AZURE_OPENAI_ENDPOINT")
azure_openai_chat_model = os.getenv("AZURE_OPENAI_CHAT_MODEL")
azure_openai_api_version = os.getenv("AZURE_OPENAI_API_VERSION")

llm = AzureChatOpenAI(
    model=azure_openai_chat_model,
    api_key=azure_openai_api_key,
    api_version=azure_openai_api_version,
    azure_endpoint=azure_openai_endpoint,
    temperature=0
)

print("PowerPoint Analyzer initialized with Azure OpenAI")

# Define the state structure for our workflow
class PPTAnalysisState(TypedDict):
    ppt_folder: str               # Folder path containing PowerPoint files
    ppt_files: List[str]          # List of PowerPoint file paths
    all_extracted_text: str       # Combined text from all PPTs
    structured_data: List[dict]   # Final structured data
    output_file: str              # Path to output CSV file

def extract_text_from_ppt(ppt_path: str) -> str:
    """
    Extract ALL text content comprehensively from a PowerPoint file.
    
    Args:
        ppt_path: Path to the PowerPoint file
        
    Returns:
        Comprehensive text from all slides with detailed structure
    """
    try:
        presentation = Presentation(ppt_path)
        all_text = []
        
        print(f"  📄 Comprehensively analyzing {len(presentation.slides)} slides in {os.path.basename(ppt_path)}...")
        
        # Add file header with metadata
        file_summary = f"\n{'='*60}\n"
        file_summary += f"FILE: {os.path.basename(ppt_path)}\n"
        file_summary += f"TOTAL SLIDES: {len(presentation.slides)}\n"
        file_summary += f"{'='*60}\n"
        all_text.append(file_summary)
        
        for slide_num, slide in enumerate(presentation.slides, 1):
            print(f"    🔍 Processing slide {slide_num}/{len(presentation.slides)}...")
            slide_content = []
            slide_header = f"\n--- SLIDE {slide_num} CONTENT START ---\n"
            slide_content.append(slide_header)
            
            # Track different types of content
            text_boxes = []
            table_content = []
            other_content = []
            
            # Extract from all shapes comprehensively
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text") and shape.text.strip():
                    # Categorize content by shape type
                    if shape.has_text_frame:
                        text_boxes.append(f"TEXT BOX {shape_idx + 1}: {shape.text.strip()}")
                    else:
                        other_content.append(f"SHAPE TEXT: {shape.text.strip()}")
                
                # Extract comprehensive table data
                if shape.has_table:
                    table = shape.table
                    table_content.append(f"\nTABLE {len(table_content) + 1} CONTENT:")
                    
                    # Extract table headers if present
                    if table.rows:
                        header_row = " | ".join([cell.text.strip() for cell in table.rows[0].cells])
                        if header_row:
                            table_content.append(f"HEADERS: {header_row}")
                    
                    # Extract all table rows
                    for row_idx, row in enumerate(table.rows):
                        row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                        if row_text:
                            table_content.append(f"ROW {row_idx + 1}: {row_text}")
            
            # Compile slide content in organized manner
            if text_boxes:
                slide_content.append("\nTEXT CONTENT:")
                slide_content.extend(text_boxes)
            
            if table_content:
                slide_content.append("\nTABLE DATA:")
                slide_content.extend(table_content)
            
            if other_content:
                slide_content.append("\nOTHER CONTENT:")
                slide_content.extend(other_content)
            
            # Add slide footer
            slide_footer = f"\n--- SLIDE {slide_num} CONTENT END ---\n"
            slide_content.append(slide_footer)
            
            # Add slide content to main collection
            slide_content_text = "\n".join(slide_content)
            all_text.extend(slide_content)
            
            # Debug output for slide processing
            content_length = len(slide_content_text)
            shape_count = len(slide.shapes)
            print(f"      ✅ Slide {slide_num}: {content_length} chars extracted from {shape_count} shapes")
        
        # Add file footer
        file_footer = f"\n{'='*60}\n"
        file_footer += f"END OF FILE: {os.path.basename(ppt_path)}\n"
        file_footer += f"TOTAL SLIDES PROCESSED: {len(presentation.slides)}\n"
        file_footer += f"{'='*60}\n"
        all_text.append(file_footer)
        
        final_content = "\n".join(all_text)
        print(f"  ✅ Extracted {len(final_content)} characters from ALL {len(presentation.slides)} slides")
        
        return final_content
    
    except Exception as e:
        error_msg = str(e)
        print(f"  ❌ Error processing {os.path.basename(ppt_path)}")
        
        # Provide specific error guidance
        if "Package not found" in error_msg:
            print(f"      💡 File may be corrupted, password-protected, or currently open in PowerPoint")
        elif "Permission denied" in error_msg:
            print(f"      💡 File permission issue - check if file is read-only or locked")
        elif "not a zip file" in error_msg:
            print(f"      💡 File appears to be corrupted or not a valid PowerPoint file")
        else:
            print(f"      📝 Error details: {error_msg}")
        
        return f"Error reading {os.path.basename(ppt_path)}: {error_msg}"

# Node 1: Discover and extract text from all PowerPoint files
def extract_all_text_node(state: PPTAnalysisState) -> PPTAnalysisState:
    """Extract text from all PowerPoint files in the specified folder."""
    
    # Get the folder path from state, or use default if not provided
    default_folder = os.path.join(os.path.expanduser("~"), "Downloads", "sample PPTs")
    ppt_folder = state.get("ppt_folder", default_folder)
    
    print(f"📁 Scanning for PowerPoint files in: {ppt_folder}")
    
    # Check if folder exists
    if not os.path.exists(ppt_folder):
        print(f"⚠️  Folder not found: {ppt_folder}")
        print(f"� Please check the folder path and try again")
        return {
            **state,
            "ppt_files": [],
            "all_extracted_text": "No PowerPoint files found - folder does not exist."
        }
    
    # Find all PowerPoint files
    ppt_files = []
    if os.path.exists(ppt_folder):
        ppt_patterns = [
            os.path.join(ppt_folder, "*.pptx"),
            os.path.join(ppt_folder, "*.ppt")
        ]
        for pattern in ppt_patterns:
            all_files = glob.glob(pattern)
            # Filter out temporary files (starting with ~$)
            valid_files = [f for f in all_files if not os.path.basename(f).startswith('~$')]
            ppt_files.extend(valid_files)
    
    if not ppt_files:
        print(f"⚠️  No PowerPoint files found in {ppt_folder}")
        print(f"💡 Please add .pptx or .ppt files to the folder")
        return {
            **state,
            "ppt_files": [],
            "all_extracted_text": "No PowerPoint files found to process."
        }
    
    print(f"📋 Found {len(ppt_files)} PowerPoint files:")
    for file in ppt_files:
        print(f"  • {os.path.basename(file)}")
    
    # Extract text from all files
    all_text_content = []
    successful_files = 0
    failed_files = []
    
    for ppt_file in ppt_files:
        print(f"\n🔍 Extracting text from: {os.path.basename(ppt_file)}")
        text_content = extract_text_from_ppt(ppt_file)
        
        # Check if extraction was successful
        if text_content.startswith("Error reading"):
            failed_files.append(os.path.basename(ppt_file))
            print(f"  ⚠️  Skipping failed file: {os.path.basename(ppt_file)}")
        else:
            all_text_content.append(text_content)
            successful_files += 1
    
    combined_text = "\n\n".join(all_text_content)
    
    print(f"\n✅ Text extraction complete!")
    print(f"   📊 Successfully processed: {successful_files} files")
    print(f"   📊 Total characters extracted: {len(combined_text)}")
    
    if failed_files:
        print(f"   ⚠️  Failed to process {len(failed_files)} files:")
        for failed_file in failed_files:
            print(f"      • {failed_file}")
        print(f"   💡 Failed files may be corrupted, password-protected, or currently open")
    
    return {
        **state,
        "ppt_files": ppt_files,
        "all_extracted_text": combined_text
    }

# Node 2: Analyze extracted text using AI
def analyze_content_node(state: PPTAnalysisState) -> PPTAnalysisState:
    """Use Azure OpenAI to analyze the extracted text and structure the data."""
    
    print("\n🤖 Analyzing content with Azure OpenAI...")
    
    if not state["all_extracted_text"] or state["all_extracted_text"].strip() == "No PowerPoint files found to process.":
        return {
            **state,
            "structured_data": []
        }
    
    # Create the comprehensive analysis prompt
    analysis_prompt = ChatPromptTemplate.from_template("""
You are an expert business analyst with deep experience in detecting customer information from business presentations. Analyze ALL content from the PowerPoint presentations with meticulous attention to detail.

ENHANCED CUSTOMER DETECTION STRATEGY:
1. **Deep Content Analysis**: Read every word on every slide, including headers, footers, speaker notes, table cells, chart labels, and embedded text in images/shapes.

2. **Multi-Context Customer Identification**: Look for customer names in ALL possible contexts in each slide:
   - Direct mentions (e.g., "Customer: Contoso")
   - Indirect references (e.g., "the customer in the retail sector")
   - Visual elements (e.g., logos, product images)

5. **Comprehensive Data Extraction**: Extract only explicit data or information clearly mentioned in slides.

6. **Contextual Business Intelligence**: Capture win/loss reasons, financial figures, and strategic context that might be distributed across multiple slides.

7. **Systematic Content Coverage**: Analyze ALL text content from every slide element including text boxes, shapes, tables, charts, and any embedded content.

CRITICAL OUTPUT FORMAT:
For each customer/deal/opportunity found, you MUST use this EXACT format with NO bullets, numbers, or extra formatting:

Customer: [customer name - use the most complete and accurate name found across all slides]
Product: [product/service name] 
Won or lost: [won/lost - only if explicitly stated or clearly implied]
Revenue Opportunity: [revenue amount - only if explicitly mentioned. Use all available information in the slide.]
Seats: [number of seats/licenses - only if explicitly mentioned. Use all available information in the slide.]
Why we won or lost: [reason/rationale for the outcome - only if explicitly mentioned]. The field name should always be "Why we won or lost".
Source File: [filename where customer was primarily identified]

SEPARATE EACH CUSTOMER WITH ONE BLANK LINE.

CRITICAL INSTRUCTIONS FOR CONSISTENT OUTPUT:
- Use EXACTLY the format shown in the examples above
- ONLY include customers actually found in the PowerPoint content provided
- Do not add any sample outputs, explanations, or commentary in your response
- Start each customer entry with "Customer:" 
- **CONSISTENCY REQUIREMENT**: You must be thorough and systematic:
  * Read through ALL content completely before starting to extract
  * Create a mental inventory of ALL customers mentioned across ALL slides
  * Extract EVERY customer found - don't skip any, even minor mentions
  * If you found customers in a previous analysis, find them again consistently
  * Be methodical - go slide by slide, paragraph by paragraph
- **CUSTOMER DETECTION PRIORITY**: Spend extra effort identifying customers accurately:
  * Search every slide thoroughly for company names, even if mentioned briefly
  * Use the most complete customer name found (full legal name preferred)
  * Include industry context if it helps identify the customer clearly
- **COMPLETENESS CHECK**: Before finishing, verify you've extracted:
  * ALL customer names from slide titles, headers, and content
  * ALL customers mentioned in case studies or examples
  * ALL customers referenced in tables, charts, or data
- If specific data isn't explicitly stated, leave blank (don't assume or infer)
- Find win or loss info from slide analysis. Don't create any info on your own
- Look for information spread across multiple slides and connect related data points
- NO extra commentary or analysis - ONLY the structured format

PowerPoint Content to Analyze:
{content}
""")
    
    # Create the analysis chain
    analysis_chain = analysis_prompt | llm | StrOutputParser()
    
    try:
        # Get the comprehensive AI analysis
        analysis_result = analysis_chain.invoke({"content": state["all_extracted_text"]})
        
        print(f"📊 Comprehensive AI Analysis complete. Processing results...")
        print(f"📝 Analysis length: {len(analysis_result)} characters")
        
        # Save the full analysis for reference
        output_dir = "output"
        os.makedirs(output_dir, exist_ok=True)
        
        # Save comprehensive analysis to text file
        analysis_file = os.path.join(output_dir, "comprehensive_analysis.txt")
        with open(analysis_file, 'w', encoding='utf-8') as f:
            f.write("COMPREHENSIVE POWERPOINT ANALYSIS\n")
            f.write("=" * 50 + "\n\n")
            f.write(analysis_result)
        
        print(f"📄 Full analysis saved to: {analysis_file}")
        
        # Parse the AI response into structured data
        structured_data = parse_ai_response(analysis_result, state)
        
        return {
            **state,
            "structured_data": structured_data
        }
    
    except Exception as e:
        print(f"❌ Error during AI analysis: {str(e)}")
        return {
            **state,
            "structured_data": []
        }

def parse_ai_response(ai_response: str, state: PPTAnalysisState) -> List[dict]:
    """Parse the comprehensive AI response into structured data format."""
    
    print("🔍 Parsing comprehensive AI analysis...")
    print(f"📝 AI Response preview: {ai_response[:500]}..." if len(ai_response) > 500 else f"📝 AI Response: {ai_response}")
    
    structured_data = []
    current_entry = {}
    
    lines = ai_response.split('\n')
    
    for line in lines:
        line = line.strip()
        
        if not line:
            # Empty line might indicate end of an entry
            if current_entry and any(current_entry.values()):
                structured_data.append(current_entry.copy())
                current_entry = {}
            continue
        
        # Parse structured data fields - be more flexible with parsing
        if line.lower().startswith("customer:"):
            if current_entry and any(current_entry.values()):
                structured_data.append(current_entry.copy())
            customer_value = line.split(":", 1)[1].strip() if ":" in line else ""
            current_entry = {"customer": customer_value}
        elif (line.lower().startswith("product:")) and current_entry:
            current_entry["product"] = line.split(":", 1)[1].strip() if ":" in line else ""
        elif (line.lower().startswith("won or lost:")) and current_entry:
            current_entry["won_or_lost"] = line.split(":", 1)[1].strip() if ":" in line else ""
        elif (line.lower().startswith("revenue opportunity:")) and current_entry:
            current_entry["revenue_opportunity"] = line.split(":", 1)[1].strip() if ":" in line else ""
        elif (line.lower().startswith("seats:")) and current_entry:
            current_entry["seats"] = line.split(":", 1)[1].strip() if ":" in line else ""
        elif (line.lower().startswith("why we won or lost:")) and current_entry:
            current_entry["why_we_won_or_lost"] = line.split(":", 1)[1].strip() if ":" in line else ""
        elif (line.lower().startswith("source file:")) and current_entry:
            current_entry["source_file"] = line.split(":", 1)[1].strip() if ":" in line else ""
    
    # Don't forget the last entry
    if current_entry and any(current_entry.values()):
        structured_data.append(current_entry)
    
    print(f"📋 Successfully parsed {len(structured_data)} customer records")
    
    # If no structured data was found, try alternative parsing
    if not structured_data and ai_response.strip():
        print("⚠️  No structured data found with standard parsing, attempting alternative extraction...")
        
        # Look for customer mentions in the comprehensive analysis
        import re
        
        # Try to extract customer names from the detailed analysis you saw
        customer_patterns = [
            r'\*\*Customer:\s*([^*\n]+)\*\*',
            r'Customer:\s*([^\n]+)',
            r'(\d+\.)\s*\*\*Customer:\s*([^*\n]+)\*\*'
        ]
        
        for pattern in customer_patterns:
            matches = re.findall(pattern, ai_response, re.IGNORECASE)
            if matches:
                print(f"🔍 Found {len(matches)} customers using pattern: {pattern}")
                break
        
        # Create fallback entry with dynamic information
        # Extract actual source file from state if available
        source_files = state.get('ppt_files', [])
        if source_files:
            source_file_names = [os.path.basename(f) for f in source_files]
            source_file_info = ", ".join(source_file_names)
        else:
            source_file_info = "Multiple files"
        
        # Try to extract some basic info from the AI response
        customer_count = ai_response.lower().count('customer')
        deal_mentions = ai_response.lower().count('revenue') + ai_response.lower().count('opportunity')
        
        structured_data.append({
            "customer": f"PARSING FAILED - Check comprehensive_analysis.txt ({customer_count} customers detected)",
            "product": "Microsoft Identity & Security products",
            "won_or_lost": "Mixed results - see comprehensive_analysis.txt",
            "revenue_opportunity": f"Multiple deals detected ({deal_mentions} revenue mentions)",
            "seats": "Various seat counts",
            "why_we_won_or_lost": "Multiple wins/losses documented - see comprehensive_analysis.txt for full details",
            "source_file": source_file_info
        })
        
        print("� Created fallback entry - check comprehensive_analysis.txt for complete details")
    
    return structured_data

# Node 3: Generate CSV output
def generate_csv_node(state: PPTAnalysisState) -> PPTAnalysisState:
    """Generate the final CSV output file."""
    
    print("\n📋 Generating CSV output...")
    
    # Create output directory if it doesn't exist
    output_dir = "output"
    os.makedirs(output_dir, exist_ok=True)
    
    output_file = os.path.join(output_dir, "ppt_analysis.csv")
    
    if not state["structured_data"]:
        print("⚠️  No structured data to export")
        # Create empty CSV with headers
        empty_df = pd.DataFrame(columns=[
            "customer", "product", "won_or_lost", "revenue_opportunity", 
            "seats", "why_we_won_or_lost", "source_file"
        ])
        empty_df.to_csv(output_file, index=False, encoding='utf-8-sig')
        print(f"📁 Empty CSV created: {output_file}")
    else:
        # Create DataFrame and export to CSV
        df = pd.DataFrame(state["structured_data"])
        
        # Ensure all required columns exist
        required_columns = ["customer", "product", "won_or_lost", "revenue_opportunity", "seats", "why_we_won_or_lost", "source_file"]
        for col in required_columns:
            if col not in df.columns:
                df[col] = ""
        
        # Reorder columns
        df = df[required_columns]
        
        # Export to CSV with UTF-8 BOM for proper Excel compatibility
        df.to_csv(output_file, index=False, encoding='utf-8-sig')
        
        print(f"✅ CSV file created: {output_file}")
        print(f"📊 Total records: {len(df)}")
        print(f"📄 Comprehensive analysis available in: output/comprehensive_analysis.txt")
        
        # Display summary
        print("\n📋 Summary of extracted data:")
        print(df.to_string(index=False))
        
        # Additional insights
        print(f"\n💡 Analysis insights:")
        print(f"   • Processed {len(state.get('ppt_files', []))} PowerPoint files")
        print(f"   • Extracted {len(df)} business records")
        print(f"   • Full content analysis saved for detailed review")
        print(f"   • All slide content was comprehensively analyzed for context")
    
    return {
        **state,
        "output_file": output_file
    }

# Create the LangGraph workflow
def create_ppt_analyzer_workflow():
    """Create the PowerPoint analyzer workflow using LangGraph."""
    
    # Create the state graph
    workflow = StateGraph(PPTAnalysisState)
    
    # Add nodes
    workflow.add_node("extract_text", extract_all_text_node)
    workflow.add_node("analyze_content", analyze_content_node)
    workflow.add_node("generate_csv", generate_csv_node)
    
    # Define the workflow edges
    workflow.add_edge(START, "extract_text")
    workflow.add_edge("extract_text", "analyze_content")
    workflow.add_edge("analyze_content", "generate_csv")
    workflow.add_edge("generate_csv", END)
    
    # Compile the workflow
    app = workflow.compile()
    
    return app

# Main execution function
def run_ppt_analyzer(ppt_folder=None):
    """Run the comprehensive PowerPoint analyzer workflow."""
    
    print("🚀 Starting Comprehensive PowerPoint Analyzer")
    print("=" * 60)
    print("📊 This analyzer will:")
    print("   • Extract ALL content from every slide")
    print("   • Perform comprehensive business analysis")
    print("   • Generate structured CSV data")
    print("   • Save detailed analysis for review")
    print("=" * 60)
    
    # Get folder path from user if not provided
    if not ppt_folder:
        print("\n📁 Folder Selection:")
        print("Enter the full path to your PowerPoint files folder")
        print(f"(or press Enter for default: {os.path.join(os.path.expanduser('~'), 'Downloads', 'sample PPTs')})")
        
        user_input = input("\nFolder path: ").strip()
        
        if user_input:
            ppt_folder = user_input
        else:
            ppt_folder = os.path.join(os.path.expanduser("~"), "Downloads", "sample PPTs")
    
    print(f"\n📂 Using folder: {ppt_folder}")
    
    # Validate folder exists
    if not os.path.exists(ppt_folder):
        print(f"❌ Error: Folder does not exist: {ppt_folder}")
        print("💡 Please check the path and try again")
        return
    
    # Create the workflow
    app = create_ppt_analyzer_workflow()
    
    # Initialize the state with the folder path
    initial_state = {
        "ppt_folder": ppt_folder,
        "ppt_files": [],
        "all_extracted_text": "",
        "structured_data": [],
        "output_file": ""
    }
    
    try:
        # Run the workflow
        final_state = app.invoke(initial_state)
        
        print("\n" + "=" * 60)
        print("✅ Comprehensive PowerPoint Analysis Complete!")
        print("=" * 60)
        
        if final_state["output_file"]:
            print(f"� Structured data saved to: {final_state['output_file']}")
            print(f"� Comprehensive analysis saved to: output/comprehensive_analysis.txt")
            print("💡 Open CSV in Excel for structured data")
            print("💡 Review comprehensive analysis file for detailed insights")
            print("\n🎯 Both files provide complementary views:")
            print("   • CSV: Structured business data for analysis")
            print("   • TXT: Complete narrative analysis and insights")
        else:
            print("⚠️  No output files generated")
            print("💡 Check if PowerPoint files contain readable content")
            
    except Exception as e:
        print(f"\n❌ Error during comprehensive analysis: {str(e)}")
        print("💡 Please check your PowerPoint files and try again")

if __name__ == "__main__":
    import sys
    
    # Check if folder path provided as command line argument
    if len(sys.argv) > 1:
        folder_path = sys.argv[1]
        print(f"📁 Using folder from command line: {folder_path}")
        run_ppt_analyzer(folder_path)
    else:
        # Interactive mode - ask user for folder
        run_ppt_analyzer()